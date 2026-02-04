import re
from hashlib import md5
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

def get_slide_shapes(slide):
    """Returns a list of (shape_id, fingerprint) for elements with text."""                
    contents = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            contents.append((shape.shape_id, fingerprint(shape.text.strip())))
        #SmartArt support is a pain to implement, leaving attempts for posterity in case python-pptx introduces support 
        if is_smart_art(shape):
            contents.append((shape.shape_id, fingerprint(get_smart_art_as_markdown(shape))))
    return contents

def fingerprint(text):
    return(md5(text.encode()).hexdigest())

def markdown_to_shape(shape, markdown_text):
    """Clarity - decide the shape type and call appropriate method to get Markdown."""
    if shape.has_text_frame:
        markdown_to_text_shape(shape,markdown_text)
    elif is_smart_art(shape):
        return markdown_to_smart_art(shape,markdown_text)


def markdown_to_smart_art(shape, translated_markdown):
    """Safely updates SmartArt nodes (including placeholders) from translated Markdown."""
    
    # 1. Access the SmartArt object safely
    try:
        # For PlaceholderGraphicFrame, it's often under .graphic_frame.smart_art
        # or directly as .smart_art
        smart_art = getattr(shape, 'smart_art', None)
        if not smart_art and hasattr(shape, 'graphic_frame'):
            smart_art = shape.graphic_frame.smart_art
    except Exception as e:
        print(f"  [ERROR] Could not access SmartArt for ID {shape.shape_id}: {e}")
        return

    if not smart_art:
        print(f"  [WARN] SmartArt property not found for ID {shape.shape_id}")
        return

    # 2. Clean the AI response
    lines = translated_markdown.strip().split('\n')
    # Remove markdown bullets and leading whitespace
    clean_texts = [re.sub(r'^\s*[\*\-]\s*', '', line) for line in lines]
            
    # 3. Get the actual nodes that have text frames
    nodes = [n for n in smart_art.nodes if n.text_frame]

    # 4. Map 1-to-1
    #for i in range(min(len(nodes), len(clean_texts))):
    #    # Option A: Simple plain text (Use this if you don't have links)
    #    # nodes[i].text_frame.text = clean_texts[i]
    #    # Option B: Call a helper to handle formatting/links (Recommended)
    #    #apply_markdown_to_text_frame(nodes[i].text_frame, clean_texts[i])

def markdown_to_text_shape(shape, markdown_text):
    """Parses Markdown and applies it to a PowerPoint shape."""
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame

    # 1. Capture original font sizes and colours
    sizes = []
    color_rgb = None
    color_theme = None

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            clean_text = run.text.strip()
            if not clean_text:
                continue
            if run.font.size:
                sizes.append(run.font.size.pt)
            # Capture color if not already found
            if not color_rgb and not color_theme:
                if run.font.color.type: # Check if color is defined
                    color_theme = run.font.color.theme_color
                    try:
                        color_rgb = run.font.color.rgb
                    except AttributeError:
                        color_rgb = None
    
    # Fallback to a standard size if no font info was found
    # Use average size -1 just in case translation is longer or to offset if a too big font has more runs.
    avg_size = ((sum(sizes) / len(sizes)) -1) if sizes else 18

    text_frame.clear() # Wipe existing content
    #Set autosize    
    #text_frame.auto_size = MSO_AUTO_SIZE.NONE
    #text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    #text_frame.word_wrap = True

    lines = markdown_text.split('\n')
    
    for i, line in enumerate(lines):
        # Create a new paragraph (except for the first one which is created by default)
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        
        # Detect bullet level (counting leading spaces)
        strip_line = line.lstrip()
        if strip_line.startswith('* '):
            p.level = (len(line) - len(strip_line)) // 2
            content = strip_line[2:]
        else:
            content = line

        # Regex to find [text](url), **bold**, and *italic*
        # This is a simplified parser; order of bold/italic matters
        parts = re.split(r'(\*\*.*?\*\*|\*.*?\*|\[.*?\]\(.*?\))', content)
        
        for part in parts:
            if not part: continue
            
            run = p.add_run()
            
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]
                run.font.bold = True
            elif part.startswith('*') and part.endswith('*'):
                run.text = part[1:-1]
                run.font.italic = True
            elif part.startswith('[') and '(' in part:
                # Basic Markdown Link: [text](url)
                link_text = re.search(r'\[(.*?)\]', part).group(1)
                link_url = re.search(r'\((.*?)\)', part).group(1)
                run.text = link_text
                run.hyperlink.address = link_url
            else:
                run.text = part
            # 5. Re-apply Size and Colour
            run.font.size = Pt(avg_size)
            if color_theme:
                run.font.color.theme_color = color_theme
            elif color_rgb:
                run.font.color.rgb = color_rgb



def get_shape_markdown(shape):
    """Clarity - decide the shape type and call appropriate method to get Markdown."""
    if shape.has_text_frame:
        return get_shape_as_markdown(shape)
    elif is_smart_art(shape):
        return get_smart_art_as_markdown(shape)
    else:
        return ""
    
def get_shape_as_markdown(shape):
    """Converts a PowerPoint shape's text and formatting to Markdown."""
    if not shape.has_text_frame:
        return ""
    
    md_paragraphs = []
    for paragraph in shape.text_frame.paragraphs:
        md_runs = ""
        for run in paragraph.runs:
            text = run.text
            if not text: continue
            
            # Apply formatting markers
            if run.font.bold:
                text = f"**{text}**"
            if run.font.italic:
                text = f"*{text}*"
            
            # Handle Hyperlinks
            try:
                if run.hyperlink and run.hyperlink.address:
                    text = f"[{text}]({run.hyperlink.address})"
            except (KeyError, AttributeError):
                # This happens if the rId is missing in the slide's .rels file
                # or if the hyperlink is malformed
                print(f"  [WARN] Broken hyperlink detected in shape {shape.shape_id} run: {text}. Skipping link.")
            md_runs += text
        
        # Handle Bullet Levels (Basic)
        prefix = ""
        if paragraph.level > 0:
            prefix = "  " * paragraph.level + "* "
        md_paragraphs.append(prefix + md_runs)        
    return "\n".join(md_paragraphs)

def get_smart_art_as_markdown(shape):
    """Converts nested SmartArt into Markdown, preserving links and basic formatting."""

    try:
        # This is the path for GraphicFrames to find their internal SmartArt
        smart_art = shape.graphic_frame.smart_art
    except AttributeError:
        # If it's a PlaceholderGraphicFrame, it might be directly on the shape or via .graphic
        smart_art = getattr(shape, 'smart_art', None)

    if not smart_art:
        print(f"  [DEBUG] Shape {shape.shape_id} identified as SmartArt, but property not found.")
        return ""

    markdown_lines = []
    for node in smart_art.nodes:
        if not node.text_frame:
            continue
            
        # Build the text for this specific node
        node_parts = []
        for paragraph in node.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text
                if not text.strip():
                    node_parts.append(text)
                    continue
                
                # Check for Hyperlinks
                hlink = run.hyperlink
                if hlink and hlink.address:
                    text = f"[{text}]({hlink.address})"
                
                # Check for Bold/Italic (Optional but helpful for AI context)
                if run.font.bold:
                    text = f"**{text}**"
                if run.font.italic:
                    text = f"_{text}_"
                
                node_parts.append(text)
            
            # Add a space between paragraphs within the same node if needed
            node_parts.append("\n")

        full_node_text = "".join(node_parts).strip()
        if full_node_text:
            indent = "  " * node.level
            markdown_lines.append(f"{indent}* {full_node_text}")

    return "\n".join(markdown_lines)


def is_smart_art(shape):
    # All attempts to implement support for SmartArt have failed
    # Giving up
    return False

    # Check 1: Explicit check for SmartArt inside the PlaceHolderGraphicFrames
    if type(shape).__name__ == 'PlaceholderGraphicFrame':
        if 'diagram' in shape._element.xml.lower():
            return True

    # Check 2: Standard check for IGX (Type 15)
    if shape.shape_type == MSO_SHAPE_TYPE.IGX_GRAPHIC:
        return True

    # Check 3: The getattr safety (might still work on some versions)
    if getattr(shape, "has_smart_art", False):
        return True
        
    return False
