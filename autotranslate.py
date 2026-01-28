from config import load_config
from genai import getAI
import os
import shutil
import argparse
import sqlite3
import hashlib
import copy
import json
import time
import sys
import re
from itertools import batched
from database import init_db
from pathlib import Path
from pptx import Presentation
from pptxhandler import get_slide_shapes, get_shape_markdown, markdown_to_shape, is_smart_art
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from classes import Course, create_course_object

def lang_code_to_text(code):
    mapping = {
        "fi": "Finnish",
        "sv": "Swedish",
        "en": "English",
        "es": "Spanish",
        "de": "German",
        "fr": "French"
    }
    # Fall back to Klingon if necessary
    return mapping.get(code.strip().lower(), "Klingon")


#############################################################################
# TEXT COLLECTION FOR TRANSLATION
#############################################################################

def collect_texts_from_pptx(file_path, languages,conn,list_of_texts):
    path = Path(file_path)
    cursor = conn.cursor()
    
    # Get original file's last modified time
    orig_mtime = path.stat().st_mtime

    for lang in languages:
        language_name = lang_code_to_text(lang)
        trans_path = path.parent / f"{path.stem}_{lang}{path.suffix}"
        ppath = str(path)[-30:]
        ptrans_path = str(trans_path)[-30:]
        
        # CHANGE DETECTION:
        # Trigger "Nuclear Option" if:
        # a) File doesn't exist
        # b) Original is newer than the translation
        needs_rebuild = (not trans_path.exists()) or (orig_mtime > trans_path.stat().st_mtime)
        
        if not needs_rebuild:
            print(f"[{lang}] {ppath} ==> {ptrans_path}. Translation is up to date. Skipping.")
            continue

        print(f"[{lang}] {ppath} ==> {ptrans_path}. Updated or translation missing. Collecting strings to translate...")
        prs = Presentation(path)

        ###############################################
        # COLLECT TEXTS TO TRANSLATE FROM ALL SLIDES
        ###############################################

        n = 1
        for slide in prs.slides:
            #print(f"Slide {n}...")
            shape_map = {s.shape_id: s for s in slide.shapes if hasattr(s, "text") or is_smart_art(s)}
            slide_data = get_slide_shapes(slide)
            if not slide_data:
                continue
                
            # 3.1 CHECK TLB
            for shape_id, fingerprint in slide_data:
                cursor.execute("SELECT EXISTS(SELECT 1 FROM tlb WHERE fingerprint=? AND lang_code=?)", (fingerprint,lang))
                exists = cursor.fetchone()[0]  
                if not exists:
                    list_of_texts[lang].append( {"id": fingerprint, "text": get_shape_markdown(shape_map[shape_id]), "translation": ""})
            #n += 1

#############################################################################
# TRANSLATION OF TEXT BATCHES
#############################################################################

def translate_texts(translation_texts, batch_size, prompt, model, conn):

    cursor = conn.cursor()
    
    for lang, texts in translation_texts.items():
        language_name = lang_code_to_text(lang)
        # Determine suitable batch size
        L = len(texts)
        if L == 0:
            continue
        if batch_size >= L:
            B = L
        else:
            B = batch_size
            while 0 < (L%B) < (B/2):
                B += 1 
        
        print(f"[{lang}] Calling AI for {L} new strings in batches of {B} to translate!")
                
        for translate_now in batched(texts, B):
            # Batch translate the missing snippets
            try:
                ai_results = model.translate(translate_now, prompt, "Finnish", language_name)
            except KeyboardInterrupt:
                print("\n[STOP] Keyboard Interrupt detected!")
                ai_results = []
                ai_results = translate_now
            except Exception as e:
                print(f"\n[ERROR] AI Translation failed unexpectedly: {e}")
                ai_results = translate_now

            # Open the SQL fix file in append mode
            with open("fix_translations.sql", "a", encoding="utf-8") as sql_file:
                #for shape_id, fingerprint, original_markdown, translated_markdown in ai_results:
                for entry in ai_results:
                    if entry["translation"] == "" or entry["translation"] is None:
                        #Translation failed for this string!
                        print(f"ERROR! AI provided no translation for string: '{entry["text"][:60]}...'")

                        # --- Generate SQL Manual translation Line ---
                        # We escape single quotes for SQL safety
                        sql_orig = entry['text'].replace("'", "''")
                        sql_file.write(f"/* INSERT MANUAL TRANSLATION!!!\n ORIG: {entry["text"]}\n*/\n")
                        sql_file.write(f"INSERT OR REPLACE INTO tlb (fingerprint, source_text, target_text, lang_code) VALUES ('{entry["id"]}', '{entry["text"]}', 'INSERT MANUAL TRANSLATION HERE', '{lang}');\n")
                    else:
                        # --- Update TLB and commit ---
                        cursor.execute(    
                            "INSERT OR REPLACE INTO tlb (fingerprint, source_text, target_text, lang_code) VALUES (?, ?, ?, ?)",
                            (entry["id"], entry["text"], entry["translation"], lang)
                        )
                        conn.commit()

                        # --- Generate SQL Fix Line ---
                        # We escape single quotes for SQL safety
                        sql_orig = entry['text'].replace("'", "''")
                        sql_trans = entry['translation'].replace("'", "''")            
                        sql_file.write(f"/* FIX AND UNCOMMENT IF NECESSARY!!!\n ORIG: {entry['text']}\n CURRENT: {entry['translation']}*/\n")
                        sql_file.write(f"/* UPDATE tlb SET target_text = '{sql_trans}' WHERE fingerprint = '{entry['id']}';*/\n")
                sql_file.close()
        print("\n[CHECKPOINT] Press ENTER to process the next translations, or Ctrl+C to stop here.")
        input(">> ")


###############################################
# CREATING TRANSLATED FILES
###############################################

def create_translated_pptx(file_path, languages,conn):
    path = Path(file_path)
    cursor = conn.cursor()
    
    # Get original file's last modified time
    orig_mtime = path.stat().st_mtime

    for lang in languages:
        language_name = lang_code_to_text(lang)
        trans_path = path.parent / f"{path.stem}_{lang}{path.suffix}"
        ppath = str(path)[-30:]
        ptrans_path = str(trans_path)[-30:]
        translation_ok = True
        
        # CHANGE DETECTION:
        # a) File doesn't exist
        # b) Original is newer than the translation
        needs_rebuild = not trans_path.exists() or (orig_mtime > trans_path.stat().st_mtime)
        
        if not needs_rebuild:
            continue

        prs = Presentation(path)

        for slide in prs.slides:
            shape_map = {s.shape_id: s for s in slide.shapes if hasattr(s, "text")}
            slide_data = get_slide_shapes(slide)
            if not slide_data:
                continue
                
            # Fetch translation from TLB
            for shape_id, fingerprint in slide_data:
                cursor.execute("SELECT target_text FROM tlb WHERE fingerprint=?", (fingerprint,))
                row = cursor.fetchone()
                #push markdown into shape
                if row:
                    markdown_to_shape(shape_map[shape_id], row[0])            
                else:
                    translation_ok = False
                    break
            if not translation_ok:
                break
        if translation_ok:
            prs.save(trans_path)
            print(f"[{lang}] Successfully updated {trans_path.name}")
        else:
            # No translation, delete the translated file
            # trans_path.unlink()
            print(f"[{lang}] Translation for {trans_path.name} has failed, please insert manual translations to TLB!")


#############################################################################
# MAIN
#############################################################################

if __name__ == "__main__":

    parser = argparse.ArgumentParser(description="PDF Publisher AutoTranslate")
    parser.add_argument("--silent", "-s", action="store_true", help="Silent mode, minimal output")
    args = parser.parse_args()
    silent = args.silent

    (config, publications) = load_config()
    if not silent:
        print("Config loaded successfully!")
    db = init_db()
    
    #Initialise genAI
    model = getAI(AI=config["gen_ai"]["AI"],api_key=config["gen_ai"]["API_KEY"],model=config["gen_ai"]["Model"],timeout=config["gen_ai"]["Request_timeout_ms"],maxperminute=config["gen_ai"]["Max_requests_per_minute"])

    #For every publication
    for pub in publications:
        if not silent:
            print(f"Tarkistetaan {config[pub]['coursename']} - käännökset: {config[pub]['translate_to']}") 
        if config[pub]['translate_to'] == "":
            continue
        languages = config[pub]['translate_to'].split(",")

        # Initialise course/publication object
        courseObject = create_course_object(config, pub)
        texts_to_translate = {}
        for lang in languages:
            texts_to_translate[lang] = []
        # Read the lecture names and topics from the configuration, error if not enough lecture definitions are found:
        try:
            for x in range(1, courseObject.lectures+1):
                lecturelist = config[pub][str(x)].split(";")
                lecture_name = lecturelist.pop(0).strip()
                courseObject.add_lecture(lecture_name, x, [topic.strip() for topic in lecturelist])
        except KeyError:
            print(f"Lectures should be added as <lecturenumber = name, topic1, topic2 ... topicN> under publication {courseObject.name} in settings.ini")
            continue

        # Determine list of files
        files = []
        files.append((Path(courseObject.course_slides_dir) / config["settings"]["headerfile"]).with_suffix(".pptx"))
        files.append((Path(courseObject.course_slides_dir) / config["settings"]["dividerfile"]).with_suffix(".pptx"))
        files.append((Path(courseObject.course_slides_dir) / config["settings"]["footerfile"]).with_suffix(".pptx"))

        # Go through lectures
        for n in range(1, courseObject.lectures+1):
            for topic in courseObject.lecture_list[n-1].topic_list:
                files.append((Path(config['settings']['lecture_slides_dir']) / topic).with_suffix(".pptx"))
            #publication-specific additional lecture slides
            files.append(next((p for p in Path(config[pub]['course_slides_dir']).glob(f"*{n:02d}*.pptx") if not re.search(r'_[a-zA-Z]{2}\.pptx$', p.name)), None))
        # Collect texts:
        for f in files:
            collect_texts_from_pptx(f, languages,db, texts_to_translate)

        # Make translations for this publication
        translate_texts(texts_to_translate,int(config['gen_ai']['batch_size']), config[pub]['ai_prompt'], model,db)

        # Make translated files:
        for f in files:
            create_translated_pptx(f,languages,db)

